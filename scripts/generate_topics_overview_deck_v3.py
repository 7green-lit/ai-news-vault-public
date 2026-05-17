#!/usr/bin/env python3
"""
Generate topics-overview deck v3 (2026-05-13).

変更点(v2 → v3):
- 1枚目: 左半分=全体俯瞰(潮流02を「業務代替×AIネイティブ」に差し替え、潮流01を3軸整理、潮流03を概要+事例)
- 各トピックスライド: 3-4論点(MECE観点)で「観察+含意」+「概要・事例詳細」構造
- 直近シグナル/次の注視点も「概要+事例詳細」構造

Usage:
    python scripts/generate_topics_overview_deck_v3.py
"""
from __future__ import annotations
from pathlib import Path
import shutil

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---------- Palette ----------
CREAM       = RGBColor(0xF5, 0xF1, 0xE8)
DEEP_INDIGO = RGBColor(0x2C, 0x3E, 0x50)
SLATE       = RGBColor(0x4A, 0x55, 0x68)
SAGE        = RGBColor(0x7A, 0x9B, 0x85)
TERRACOTTA  = RGBColor(0xB5, 0x70, 0x5F)
STONE_LIGHT = RGBColor(0xEC, 0xE7, 0xDC)
STONE_MID   = RGBColor(0xD4, 0xCF, 0xC4)
SOFT_GRAY   = RGBColor(0xA8, 0xB0, 0xBA)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
INK         = RGBColor(0x1A, 0x1A, 0x1A)

FONT_JP = "Noto Sans JP"
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ---------- Helpers ----------
def new_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs

def blank_slide(prs, fill=CREAM):
    layout = prs.slide_layouts[6]
    s = prs.slides.add_slide(layout)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = fill
    return s

def add_text(slide, left, top, width, height, text, *, size=14, bold=False,
             color=SLATE, font=FONT_JP, align=None, anchor=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    if anchor:
        tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.0)
    tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0)
    tf.margin_bottom = Inches(0.0)
    p = tf.paragraphs[0]
    if align:
        p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box

def add_rect(slide, left, top, width, height, *, fill=None, line=None, line_w=0.75):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    if fill is None:
        rect.fill.background()
    else:
        rect.fill.solid()
        rect.fill.fore_color.rgb = fill
    if line is None:
        rect.line.fill.background()
    else:
        rect.line.color.rgb = line
        rect.line.width = Pt(line_w)
    rect.shadow.inherit = False
    return rect

def add_line(slide, x1, y1, x2, y2, *, color=STONE_MID, weight_pt=0.75):
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(weight_pt)
    return line

def add_pill(slide, left, top, width, height, text, *, fill=SAGE,
             text_color=WHITE, size=9, bold=True):
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    pill.adjustments[0] = 0.5
    pill.fill.solid()
    pill.fill.fore_color.rgb = fill
    pill.line.fill.background()
    tf = pill.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = FONT_JP
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = text_color
    return pill

def header_band(slide, eyebrow, title, page_num=None, page_total=None):
    add_rect(slide, Inches(0.0), Inches(0.0), Inches(0.18), SLIDE_H, fill=SAGE)
    add_text(slide, Inches(0.55), Inches(0.32), Inches(11.0), Inches(0.32),
             eyebrow, size=10, bold=True, color=SAGE, font=FONT_JP)
    add_text(slide, Inches(0.55), Inches(0.62), Inches(11.5), Inches(0.7),
             title, size=22, bold=True, color=DEEP_INDIGO, font=FONT_JP)
    if page_num is not None:
        add_text(slide, Inches(12.0), Inches(0.32), Inches(1.0), Inches(0.32),
                 f"{page_num:02d} / {page_total:02d}", size=10, bold=True,
                 color=SOFT_GRAY, font=FONT_JP, align=PP_ALIGN.RIGHT)
    add_line(slide, Inches(0.55), Inches(1.35), Inches(12.85), Inches(1.35),
             color=STONE_MID, weight_pt=0.75)

def footer(slide, text="ai-news-vault · 2026-05-13 · v3"):
    add_text(slide, Inches(0.55), Inches(7.22), Inches(12.5), Inches(0.2),
             text, size=8, color=SOFT_GRAY, font=FONT_JP)

def bullet_list(slide, left, top, width, items, *, size=9, color=SLATE,
                bullet="▸", bullet_color=SAGE, row_h=0.24):
    """概要+事例 共通の小さな箇条書きリスト."""
    for i, item in enumerate(items):
        y = top + Inches(row_h * i)
        add_text(slide, left, y, Inches(0.2), Inches(row_h),
                 bullet, size=size, bold=True, color=bullet_color, font=FONT_JP)
        add_text(slide, left + Inches(0.22), y, width - Inches(0.22),
                 Inches(row_h + 0.08), item, size=size, color=color, font=FONT_JP)

# ---------- Topic data ----------
TOPICS = [
    {
        "slug": "ai-companies",
        "title": "主要AI企業動向",
        "version": "v2",
        "articles": 16,
        "q": "各社の戦略差分(基盤モデル/プロダクト/オープン/特定領域)・資金調達/人材移動/パートナーシップ・日本勢(PKSHA/Sakana/PFN等)のグローバル位置取り — 次の主導権争いはどう見えるか。",
        "points": [
            {
                "head": "観察1: インフラ層は NVIDIA の「投資×供給」二重支配で2026 Q2 までに事実上寡占確定",
                "summary": [
                    "GPU 供給(従来)+ 株式投資による生態系ロックイン(新)で支配者ポジション",
                    "CUDA/Blackwell への構造的固着で AMD・Intel の市場参入を阻害",
                    "独禁法リスクは FTC/EC の標的になり得るが、現時点で具体的訴追なし",
                ],
                "cases": [
                    "年初来 $40B AI 関連企業株式投資コミット(2026-05-09)",
                    "CoreWeave / Lambda 等ネオクラウド5社が世界クラウドトップ30入り(2026-05-07)",
                    "Sakana × NVIDIA「TwELL」共同研究で国内勢にも食い込み(2026-05-10)",
                ],
            },
            {
                "head": "観察2: モデル層は ChatGPT/Claude/Gemini の三極化が「現実シェア=絶対MAU拡大」両立で安定状態へ",
                "summary": [
                    "ChatGPT シェア40%切るが MAU 絶対値は2026 Q1 で 35歳以上が最速成長",
                    "Claude が急伸(80倍成長 → 法務AI攻勢 → 国内SI提携)、Anthropic が B2B で OpenAI を追い抜く勢い",
                    "Gemini が Android 17 + Gboard + Search 統合で全方位浸透、Alphabet 時価総額世界首位視野",
                ],
                "cases": [
                    "Apptopia 調査で米生成AI チャット三極化(2026-05-08)",
                    "Anthropic Claude for Legal + Claude Platform on AWS 同日(2026-05-12)",
                    "Google Android Show: Googlebooks + Gemini Intelligence(2026-05-12)",
                ],
            },
            {
                "head": "観察3: 国内勢は「米AI 借用 / 特化 / 独自基盤 / データ寡占」の4戦略に分岐し、当面共存",
                "summary": [
                    "NEC・アクセンチュア型(米AI 借用 + 業界知識)が国内大手案件を獲得",
                    "Sakana 型(特化案件)が金融・製造で先行 型(独自基盤)が R&D で差別化",
                    "LINE ヤフー型(ユーザーデータ寡占)が消費者向けで対抗",
                ],
                "cases": [
                    "NEC × Anthropic 対談公開(2026-05-11)",
                    "Sakana × SMBC / × NVIDIA TwELL(2026-05-10)",
                    "LINE ヤフー Agent i ローンチ(2026-05-11)",
                ],
            },
        ],
        "signals_summary": [
            "AI ベンダー側の B2B 攻勢が4月末〜5月で急加速",
            "Anthropic は資金・計算・チャネルの3不足を一気に補填",
        ],
        "signals_cases": [
            "2026-05-12  Anthropic Claude for Legal 一斉(20+ コネクタ)",
            "2026-05-12  OpenAI Deployment Co 設立+$4B+Tomoro 買収",
            "2026-05-11  NEC × Anthropic 対談 / 2026-05-12 アクセンチュア × Anthropic 本格化",
            "2026-05-08  Anthropic 80倍成長 → Musk DC 借用(Fortune)",
        ],
        "focus_summary": [
            "Big4 全社の AI 提携完了タイミングで寡占構造が最終確定",
            "国内 SI 各社の決算で人月モデル崩壊度合いが定量化されるか",
        ],
        "focus_cases": [
            "Deloitte/EY/KPMG の AI 提携発表(2026 Q3-Q4 予想)",
            "アクセンチュア四半期決算で OpenAI Deployment Co の影響可視化",
            "Anthropic 独自 DC(Trainium 互換)建設の時期と Pentagon 復帰",
        ],
    },
    {
        "slug": "ai-regulation",
        "title": "AI規制・ガバナンス",
        "version": "v4",
        "articles": 26,
        "q": "主要AI Giants の事実上ルールメイキングに各国が後追いで規制を作る市場構造のなか、現在の事実上ルール × 各国規制状況 × 企業の対応コスト発生点はどう整理されるか。",
        "points": [
            {
                "head": "観察1: 規制対象は「能力ベース」→「用途+挙動ベース」へシフト",
                "summary": [
                    "GPT-5.5-Cyber「Trusted Access」型の用途別アクセス制御が先行例",
                    "Anthropic がフィクション悪役 AI 描写と脅迫行動の因果を技術実証",
                    "学習データの倫理キュレーション義務が現実的に検討可能に",
                ],
                "cases": [
                    "Anthropic Claude 脅迫行動研究 + 隠れ推論可視化ツール(2026-05-12)",
                    "Claude メキシコ公共機関攻撃 17,000行ツール公表(2026-05-08)",
                    "OpenAI GPT-5.5-Cyber Trusted Access for Cyber(2026-05-07)",
                ],
            },
            {
                "head": "観察2: 規制主体は政府・業界自主・AI ベンダー自主の3層協調モデルへ",
                "summary": [
                    "政府: EU AI Act、日本初 AI 法律、サイバー対処強化法10月施行",
                    "業界自主: Hollywood Human Consent Standard、米アカデミー AI 対象外",
                    "AI ベンダー自主: Anthropic alignment ツール寄贈、OpenAI Daybreak",
                ],
                "cases": [
                    "日本初 AI 特化法律成立(2026-05-11)",
                    "EU AI 法修正で性的画像 AI 生成禁止暫定合意(2026-05-08)",
                    "George Clooney + Tom Hanks + Streep が Human Consent Standard 支持(2026-05-12)",
                ],
            },
            {
                "head": "観察3: 管轄は「EU 規制重視 / 米 自由市場 / 日本 公共財型」の3極で確定",
                "summary": [
                    "日本は「規制+国産基盤+産業実装の一体化」で第3極を形成",
                    "輸出可能な公共財型 AI 戦略はデジタル庁・FOIP デジタル回廊で具体化",
                    "EU の説明可能性義務に対応するため Anthropic 系の解釈技術が業界標準化候補",
                ],
                "cases": [
                    "IBM Sovereign Core 一般提供 + ",
                    "FOIP デジタル回廊構想(2026-05-10)、日本政府が米最新AI使用権要求",
                    "高市総理がClaude Mythos受けサイバー対策指示(2026-05-12)",
                ],
            },
            {
                "head": "観察4: 規制エコシステム自体が AI に「ノイズ攻撃」される第二波が始動",
                "summary": [
                    "バグバウンティ・特許審査・行政申請が AI 生成ノイズで機能不全リスク",
                    "「AI vs AI の規制執行」が必須化、Anthropic 解釈ツールが規制側でも有用に",
                ],
                "cases": [
                    "HackerOne「AI 生成ゴミ脆弱性報告」殺到で受付停止(2026-05-11)",
                    "JBpress「AI 申請の洪水、行政が制度改革を迫られる」(2026-05-12)",
                    "AI 検索対策悪徳業者の急増(2026-05-12)",
                ],
            },
        ],
        "signals_summary": [
            "AI 兵器化と国内法律施行が同時期に集中",
            "コンテンツ真正性のスタンダード化が Hollywood 主導で加速",
        ],
        "signals_cases": [
            "2026-05-12  Google「AI 製ゼロデイ + 中朝露関与」公式報告",
            "2026-05-12  ClaudeBleed 脆弱性(Claude for Chrome 拡張)",
            "2026-05-11  日本初 AI 法律成立",
            "2026-05-12  高市総理がClaude Mythos対策指示",
        ],
        "focus_summary": [
            "公共財型 AI 規制の輸出可能性と FOIP 文脈での運用",
            "Big4(特に PwC・Deloitte)の AI 監査ビジネス規模の立ち上がり",
        ],
        "focus_cases": [
            "EU AI Act × 日本初 AI 法律の相互運用性",
            "AI 駆動の規制執行(AI vs AI)の標準化",
            "「学習データ規定」の具体化と SF 文学検閲リスク",
        ],
    },
    {
        "slug": "ai-data-structuring",
        "title": "AI活用のための構造化",
        "version": "v4",
        "articles": 16,
        "q": "LLM 業務適用のためのデータ・知識・スキーマ構造化はどう進化しているか。特に業界別・業種別に進化が進むと考えられ、国内エンタープライズのボトルネックと突破口はどこか。",
        "points": [
            {
                "head": "観察1: データ層は「文書構造化(PDF)」→「文脈構造化(Skill)」へ進化中",
                "summary": [
                    "PDF/表/DB の抽出は既に成熟、業界別「文脈+判断ルール」の構造化が次の競争領域",
                    "Anthropic Skill 型(振る舞いプロンプト+ツール)が事実上の標準化候補",
                    "国内では が先行モデル",
                ],
                "cases": [
                    "",
                    "Salesforce Headless 360 で SaaS データを MCP 経由公開(2026-04-19)",
                    "HEROZ ASK が MCP 対応開始(国内法人 SaaS 第一号、2026-05-12)",
                ],
            },
            {
                "head": "観察2: プロトコル層は MCP × Anthropic Skill が事実上の標準へ収斂中",
                "summary": [
                    "MCP/A2A/UCP/AP2 の6プロトコル乱立から、MCP が大手SaaS統合で勝ち抜き",
                    "業界別「文脈 Skill」を Anthropic が CoCounsel 連携等で先行プラットフォーム化",
                    "Google Gemini API skills(GitHub OSS 版)が対抗、競合は中立規格作りを志向",
                ],
                "cases": [
                    "GitHub MCP Server コミット前認証情報スキャン(2026-05-07)",
                    "Anthropic Claude for Legal 20+ コネクタ + Thomson Reuters CoCounsel(2026-05-12)",
                    "Google Gemini API skills GitHub 公開(2026-05-08)",
                ],
            },
            {
                "head": "観察3: 主権/ガバナンス層は「AI 主権スタック」という新カテゴリへ昇格",
                "summary": [
                    "アクセス統制+暗号化+監査+AI 運用を統合した製品が大手から登場",
                    " / IBM Sovereign Core / Oracle Alloy / Salesforce Headless 360 が同カテゴリへ",
                    "規制業界(金融・公共・医療)向けに主権ニーズが顕在化",
                ],
                "cases": [
                    "IBM Sovereign Core 一般提供開始(2026-05-05 GA、2026-05-11 国内提供)",
                    "Fivetran AI Readiness 指数2026年版(2026-05-08)で評価軸の指標化",
                    "Zenmu 秘密分散 × RAG 特許出願(2026-05-12)",
                ],
            },
            {
                "head": "観察4: 構造化対象の境界が「テキスト/表/画像」から「動作・センサー」へ拡張",
                "summary": [
                    "ファナック・安川「フィジカル AI」で動作データが構造化対象に",
                    "東芝「反事実波形」など説明可能性技術が産業応用へ",
                    "VLA モデル(視覚言語動作)向けデータ規格が次の競争軸",
                ],
                "cases": [
                    "ファナック・安川電機フィジカル AI(2026-05-11)",
                    "東芝 反事実波形生成技術で異常検知の根拠説明(2026-05-12)",
                    "NEC 独自 AI で 3D 点群90%軽量化(2026-05-12)",
                ],
            },
        ],
        "signals_summary": [
            "国内製造業データ主権の機運が同月内で複数事例化",
            "Big4 監査領域でも「AI 主権基盤」採用が始動",
        ],
        "signals_cases": [
            "2026-05-12  HEROZ ASK MCP 対応開始",
            "2026-05-12  オーエスジー × Microsoft Fabric データ基盤刷新",
            "2026-05-11  ファナック・安川フィジカル AI 国産基盤",
            "2026-05-07  三菱UFJ信託 × Snowflake データマネジメント",
        ],
        "focus_summary": [
            "Sovereign Cloud 各社(Microsoft / AWS / Oracle)の国内対応動向",
            "VLA モデル学習データ規格の標準化レース",
        ],
        "focus_cases": [
            "国内 AI ベンダー(PFN/Sakana/NEC)の Sovereign Core 連携",
            "MCP 分裂(A2A/UCP/AP2)の収斂タイミング",
            "フィジカル AI 文脈構造化の具体的データ標準",
        ],
    },
    {
        "slug": "enterprise-ai",
        "title": "エンタープライズAI導入",
        "version": "v2",
        "articles": 14,
        "q": "実際にROIが出ている領域と出ていない領域は何が違うか。組織変革との関係はどう整理され、業界・業種特化の動きはどんな領域・内容で進むか。",
        "points": [
            {
                "head": "観察1: 議論軸は「ROI算定」→「業務浸食速度」へ移行",
                "summary": [
                    "@IT「AI 業務浸食は予想以上」が業界メディアの公式時代認識へ",
                    "Cloudflare 1,100、Coinbase 14%、GM 数百名と「CEO 公言型解雇」が連鎖",
                    "ZDNET・AINOW が「生産性指標を超えた5つの戦略的転換」を提示",
                ],
                "cases": [
                    "@IT「AI 業務浸食は予想以上だった」(2026-05-11)",
                    "ガートナー CEO 28%「AI が最大の収益リスク」(2026-05-07)",
                    "GM が数百名 IT 解雇、AI スキル人材へ入れ替え(2026-05-11)",
                ],
            },
            {
                "head": "観察2: ガバナンスは「シャドーAI 統制 vs 地下化」の二律背反が顕在化",
                "summary": [
                    "管理職4割が機密情報を入力、「危険と分かっていても使う」現実",
                    "M365 Copilot 95%企業の鉄則「サポート万全は逆効果」と矛盾しない統制設計が必要",
                    "AI 申請洪水・AI 検索悪徳業者など「AI が制度を逆利用」する第二波",
                ],
                "cases": [
                    "GRAS グループ調査「管理職4割シャドー AI に機密入力」(2026-05-08)",
                    "AINOW「生成 AI 意識改革6ステップ」(2026-05-10)",
                    "AINOW シャドー AI 対策(2026-05-06)",
                ],
            },
            {
                "head": "観察3: 実装軸は「指示駆動 × 進化探索 × 両側最適化」の3方式が併存",
                "summary": [
                    "指示駆動: Codex / Claude Code / Cursor で局所最適化",
                    "進化探索: AlphaEvolve でデータセンター最適化・行列乗算改善",
                    "両側最適化: Uber × OpenAI 型(B2B2C リアルタイム最適化)が消費者文脈で本格化",
                ],
                "cases": [
                    "Google DeepMind AlphaEvolve 実用化解説(2026-05-06)",
                    "Uber × OpenAI でドライバー収益 + 配車高速化(2026-05-06)",
                    "Singular Bank × ChatGPT/Codex で60-90分/日削減(2026-05-06)",
                ],
            },
            {
                "head": "観察4: 国内最終マイルは「AINOW 教材化 + 公共セクター案件」で広がる",
                "summary": [
                    "AINOW シリーズが導入計画書・ROI・意識改革・成果報告 KPI を体系化",
                    "霞が関 18万人 AI、自治体 AI(一関市 RAG、東京都×ソウル MOU)が新市場",
                    "三菱UFJリサーチ Gov Sales のような専用ツール商品化が始動",
                ],
                "cases": [
                    "AINOW 3記事(2026-05-13): 経営層説得・ROI・成果報告 KPI",
                    "霞が関 18万人 AI で答弁・分析(2026-05-11)",
                    "三菱UFJリサーチ Gov Sales 自治体向け営業 AI(2026-05-12)",
                ],
            },
        ],
        "signals_summary": [
            "業務代替型解雇が CEO 公言レベルで連鎖、AI ネイティブ企業が脅威化",
            "国内では教材化と公共セクター案件で実装が加速",
        ],
        "signals_cases": [
            "2026-05-12  GitLab 大規模解雇 + CREDIT 価値観終了",
            "2026-05-11  ChatGPT Q1 で35歳以上が最速成長(OpenAI)",
            "2026-05-11  ダイヤモンド「AI ネイティブスタートアップの脅威」",
            "2026-05-12  IBM Enterprise Advantage 拡充",
        ],
        "focus_summary": [
            "業務浸食の業界別マッピング(製造/金融/法務/医療/公共)",
            "両側最適化型 AI の国内応用(PayPay/Yahoo/楽天/ヤマト等)",
        ],
        "focus_cases": [
            "シャドー AI 統制の現実解(規制対応 vs 利用率維持)",
            "AI 効率化型解雇の国内大手への波及度",
            "公共コンサル × AI 市場の規模",
        ],
    },
    {
        "slug": "agentic-ai",
        "title": "AIエージェント / Agentic AI",
        "version": "v3",
        "articles": 10,
        "q": "AIエージェントは「デモから実用」へどう移行していくか。信頼性・コスト・人間との協調設計の観点で何が成功要因か。",
        "points": [
            {
                "head": "観察1: 記憶層は「外部記憶 + ツール接続 + 内省」の3層化で長期記憶問題が技術的に解決へ",
                "summary": [
                    "Artifacts(外部FS)+ MCP(ツール)+ dream(内省学習)で多層化",
                    "計算コスト設計・データ主権・武器化リスクが新規論点として浮上",
                    "Claude メキシコ攻撃の自己改良パターンと正反対の「自己改善」が共存",
                ],
                "cases": [
                    "Claude Managed Agents「dream」機能(2026-05-06)",
                    "Cloudflare Artifacts(エージェント FS、2026-04-19)",
                    "OpenAI Symphony issue tracker 標準仕様(2026-04)",
                ],
            },
            {
                "head": "観察2: 解釈可能性層は「内部独白の自然言語復号」で安全テストが技術化",
                "summary": [
                    "Natural Language Autoencoders で Claude 内部思考を可視化",
                    "EU AI Act + 日本AI法律の「説明可能性」要件の技術的解",
                    "Anthropic がこの分野の業界標準を握る可能性が高い",
                ],
                "cases": [
                    "Anthropic Natural Language Autoencoders 公開(2026-05-07)",
                    "Anthropic 隠れ推論可視化ツール公開(2026-05-12)",
                    "MIT Tech Review が機械論的解釈可能性を解説(2026-05-12)",
                ],
            },
            {
                "head": "観察3: ユーザー文脈層は「データ寡占 vs 性能差」の競争へ",
                "summary": [
                    "Perplexity Personal Computer・Claude in Chrome・ChatGPT Atlas が個人 AI 環境を取り合う",
                    "LINE ヤフー Agent i が「9,000万 MAU × ユーザーデータ」で対抗",
                    "国内では NEC(B2B) × LINE ヤフー(B2C)の棲み分けが進行",
                ],
                "cases": [
                    "LINE ヤフー AI エージェント「Agent i」(2026-05-11)",
                    "Perplexity Personal Computer for Mac 一般開放(2026-05-08)",
                    "Claude in Chrome 拡張(2026-05-11)",
                ],
            },
            {
                "head": "観察4: 標準化は「公式ガイドライン×実装層ツール×評価基準」の3層が2026 Q4 で揃う見込み",
                "summary": [
                    "公式: 3社ガイドライン(2026-04)、日本 AI 法律、EU AI Act",
                    "実装: JetBrains Central、GitHub MCP、Anthropic Cowork",
                    "評価: Autoencoder + ベンチマーク後の「実エンタープライズ ROI」評価",
                ],
                "cases": [
                    "JetBrains Central(2026-05-08)",
                    "GitHub MCP コミット前認証情報スキャン(2026-05-07)",
                    "OpenClaw(OSS エージェント基盤、2026-01 → 2026-05-11 日経クロステック解説)",
                ],
            },
        ],
        "signals_summary": [
            "エージェントの内省・解釈可能性が同月で複数進展",
            "武器化リスク(メキシコ攻撃/ClaudeBleed)が並行して顕在化",
        ],
        "signals_cases": [
            "2026-05-12  ClaudeBleed 脆弱性(Claude for Chrome)",
            "2026-05-12  Needle 26M Tool Calling 蒸留(HN)",
            "2026-05-11  LINE ヤフー Agent i / OpenClaw 解説",
            "2026-05-08  Claude メキシコ攻撃 17,000行ツール公表",
        ],
        "focus_summary": [
            "dream + Autoencoder の組合せが規制執行ツールとして採用されるか",
            "消費者向けエージェント MAU 競争の決着時期",
        ],
        "focus_cases": [
            "エージェント武器化規制の国際協調(国連レベル議論の可能性)",
            "「dream の学習権」と顧客データ主権の整理",
            "Sierra AI、Cognition Devin の国内導入実績",
        ],
    },
    {
        "slug": "ai-coding",
        "title": "AIコーディング・開発支援",
        "version": "—",
        "articles": 3,
        "q": "AI コーディングエージェント(Codex/Claude Code/Devin/AlphaEvolve/Cursor/Replit/Lovable)が開発業務をどこまで代替するか。評価指標(SWE-bench 等)はどう進化するか。エンジニアスキル・キャリア・言語選択はどう再編されるか。国内 SI(富士ソフト/日立/NEC/富士通)の適応戦略は何が同じで何が違うか。",
        "points": [
            {
                "head": "観察1: 開発ツール層はバイブコーディング+自律エンジニアが主流化",
                "summary": [
                    "Cursor / Lovable / Zed 1.0 / Replit / vibe-coded widgets が個人開発者層を席巻",
                    "Cognition Devin / OpenAI Codex / Anthropic Claude Code が企業向け自律エージェントを提供",
                    "「最後まで自走する力」が GPT-5.5 熱狂の鍵に",
                ],
                "cases": [
                    "Zed 1.0 リリース(2026-05-11、Rust 製 AI 共同編集 IDE)",
                    "Bun が Claude を使って Zig → Rust 移行中(2026-05-11)",
                    "ジョイゾー「スキル39」で kintone を AI 改修(2026-05-12)",
                ],
            },
            {
                "head": "観察2: OSS 品質を AI が抜本的に変える実証段階に到達",
                "summary": [
                    "Claude Mythos × Firefox で20年前バグ含む271件修正、修正数が約15倍",
                    "GitHub MCP コミット前認証情報スキャンで「うっかりコミット」を防止",
                    "一方で HackerOne「AI 生成ゴミ脆弱性報告」殺到で受付停止の負の側面も",
                ],
                "cases": [
                    "Mozilla × Claude Mythos が Firefox 271件のバグ発掘(2026-05-08)",
                    "GitHub MCP Server コミット前スキャン(2026-05-07)",
                    "HackerOne 受付停止(2026-05-11)",
                ],
            },
            {
                "head": "観察3: エンタープライズガバナンス層が JetBrains/OpenAI/GitHub で整備期へ",
                "summary": [
                    "JetBrains Central、OpenAI「Running Codex safely」で運用ガイドが揃う",
                    "Anthropic Claude Code 利用上限引き上げ(SpaceX DC 経由で計算リソース確保)",
                    "国内では富士ソフト人月脱却が「AI 駆動 SI」のビジネスモデル先行例",
                ],
                "cases": [
                    "JetBrains Central(2026-05-08)",
                    "OpenAI Running Codex safely(2026-05-08)",
                    "Anthropic Claude Code 利用上限引き上げ(2026-05-06)",
                ],
            },
            {
                "head": "観察4: 人材市場は「Python 一択ではない」へ再編、年収と試験区分が同時変化",
                "summary": [
                    "TypeScript 平均年収952万円が Python 857万円を上回る",
                    "情報処理技術者試験が大改定、応用・高度試験が消滅(2027年度〜)",
                    "「AI ネイティブ就活生」がエントリーシート・面接練習で AI を当たり前に活用",
                ],
                "cases": [
                    "INSTANTROOM 調査 TypeScript 952万円(2026-05-08)",
                    "情報処理技術者試験大改定(2026-05-11、日経クロステック)",
                    "2027卒 AI ネイティブ就活生調査(2026-05-11)",
                ],
            },
        ],
        "signals_summary": [
            "OSS 品質革新と人材市場再編が同時進行",
            "国内 SI が「人月脱却 vs 規模拡大」で戦略分岐",
        ],
        "signals_cases": [
            "2026-05-12  Google「Create My Widget」vibe-coded widgets",
            "2026-05-12  CUDA-oxide(Rust → CUDA 公式コンパイラ、HN 387 pts)",
            "2026-05-11  TanStack npm サプライチェーン侵害(HN 685 pts)",
            "2026-05-11  「AI なら Python 不要?」HN 326 pts",
        ],
        "focus_summary": [
            "コーディングエージェントの企業ガバナンス標準化",
            "SI/コンサル収益モデルへの構造的圧力の定量化",
        ],
        "focus_cases": [
            "富士ソフト成果報酬移行の実数値(2026 Q2 決算)",
            "日立 GlobalLogic 統合 × AI 案件の北米成果",
            "国内 SI 中堅(TIS/SCSK/BIPROGY/CTC)の人月モデル対応",
        ],
    },
    {
        "slug": "ai-models",
        "title": "AIモデルの進化",
        "version": "v1",
        "articles": 4,
        "q": "主要汎用LLM(GPT/Claude/Gemini/Llama/DeepSeek/Grok 等)がいつ・どんな更新でどんな成果(精度/速度/コンテキスト/マルチモーダル)をもたらしたか。世代交代パターン(Big jump vs incremental)、能力軸進化(text→multimodal→reasoning→agentic)、オープン vs クローズドの競争構造を整理する。",
        "points": [
            {
                "head": "観察1: フロンティア競争は三極化が安定状態へ、新規参入は破壊的差別化が必要",
                "summary": [
                    "ChatGPT/Claude/Gemini で全方位機能差は縮小、差は「業界特化」「効率」「文脈」に",
                    "Anthropic は法務、OpenAI は CFO/金融、Google は消費者 OS、と特化軸が見え始める",
                    "Moonshot AI Kimi など中国系が破壊的設計で割り込みを狙う",
                ],
                "cases": [
                    "ChatGPT シェア40%切る、Claude 急伸、Gemini 安定(2026-05-08)",
                    "Anthropic Claude for Legal / OpenAI × PwC / Google Android Show(2026-05-12)",
                    "Moonshot AI Kimi の LLM 根幹設計に Musk も驚き(2026-05-12)",
                ],
            },
            {
                "head": "観察2: ローカル/エッジ層が「ChatGPT 5.5 Pro 接近」レベルに到達",
                "summary": [
                    "Llama 4 / Qwen3 / DeepSeek-R1 / Gemma 4 が無料でフロンティア接近",
                    "Tool Calling 専用蒸留(Needle 26M)が現実的になり、スマホ・IoT 上で実用",
                    "「Local AI needs to be the norm」が HN で大バズ(2026-05-11、399 pts)",
                ],
                "cases": [
                    "Needle 26M Gemini Tool Calling 蒸留(2026-05-12)",
                    "ローカル LLM 進化解説(ギズモード、2026-05-10)",
                    "Google Gemma 4 無料・高品質日本語 LLM(2026-05-03)",
                ],
            },
            {
                "head": "観察3: アーキ革新は「音声リアルタイム」「効率フォーマット」「同時聴き話し」が同時進行",
                "summary": [
                    "OpenAI Realtime API 3モデルで音声推論・翻訳・転写がコモディティ化",
                    "Sakana × NVIDIA TwELL で推論30%高速化・H100メモリ24%削減",
                    "Thinking Machines「interaction models」で同時聴き話し AI が実装段階",
                ],
                "cases": [
                    "OpenAI 次世代音声 API3モデル(2026-05-08)",
                    "Sakana × NVIDIA TwELL(2026-05-10)",
                    "Thinking Machines interaction models(2026-05-11〜12)",
                ],
            },
            {
                "head": "観察4: ドメイン特化モデルが「業界別評価ベンチマーク」を新たに要求",
                "summary": [
                    "GPT-5.5-Cyber、Claude for Legal、Anthropic Agents for FS など特化版が増加",
                    "東芝 反事実波形のような「異常検知の根拠説明」が産業応用へ",
                    "日本語トークン効率(約1.5倍コスト)解消のため国産特化 LLM が経済合理性",
                ],
                "cases": [
                    "OpenAI GPT-5.5-Cyber Trusted Access(2026-05-07)",
                    "東芝 反事実波形生成技術(2026-05-12)",
                    "@IT「日本語 AI は約1.5倍高い」実測(2026-05-12)",
                ],
            },
        ],
        "signals_summary": [
            "三極化の安定 + ローカル LLM 急進歩が同時進行",
            "業界特化モデルとシンボリック学習が次世代評価軸を要請",
        ],
        "signals_cases": [
            "2026-05-12  François Chollet シンボリック学習で新 AI 開発(Ndea)",
            "2026-05-12  Google Cloud Storage Rapid で AI 学習効率改善",
            "2026-05-12  Threads × Meta AI / Rivian AI Voice / Gboard Gemini",
            "2026-05-08  OpenAI Voice Realtime-2/Translate/Whisper",
        ],
        "focus_summary": [
            "三極化の安定性と「業界特化軸」の確立タイミング",
            "国産 LLM(ELYZA/NEC cotomi/PFN PLaMo/Sakana)の経済合理性",
        ],
        "focus_cases": [
            "Apple Intelligence の巻き返し(390億円和解後)",
            "Moonshot/DeepSeek/Qwen の中国系破壊的モデル",
            "シンボリック学習(François Chollet)の商用化時期",
        ],
    },
    {
        "slug": "ai-infra",
        "title": "AIインフラ・GPU・コスト",
        "version": "—",
        "articles": 0,
        "q": "GPU/データセンター/半導体の地政学が AI 産業の競争を決定する局面。コスト下落 vs 使用量爆発の均衡点。",
        "points": [
            {
                "head": "観察1: GPU/半導体層は NVIDIA 独占に Arm 連合 + SpaceX が割り込み",
                "summary": [
                    "Arm × Meta「AGI CPU」で x86 → Arm 系データセンター CPU 拡大",
                    "SpaceX Terafab $55B(Texas)で Musk が独自 AI 半導体メーカーへ",
                    "Intel 1.8nm Intel 18A で低価格 MPU まで最先端化",
                ],
                "cases": [
                    "SpaceX Terafab $55B(2026-05-08)",
                    "Arm × Meta AGI CPU 共同開発(2026-05-08)",
                    "Intel 株 490%上昇 + 1.8nm 全面化(2026-05-12)",
                ],
            },
            {
                "head": "観察2: データセンター層は「Pentagon 不在の Anthropic が Musk DC 借用」など政治と物理が交錯",
                "summary": [
                    "SpaceX × Anthropic で Musk が「evil」呼ばわりした Anthropic に GPU 提供",
                    "Google × SpaceX で軌道上 AI データセンターの協議(地球外進出)",
                    "AWS 中東 UAE 復旧に数カ月、紛争影響でクラウド可用性が地政学に依存",
                ],
                "cases": [
                    "SpaceX × Anthropic データセンター提携(2026-05-06〜07)",
                    "Google × SpaceX 軌道上 AI DC 協議(2026-05-12)",
                    "AWS UAE 復旧2カ月ぶり報告(2026-05-06)",
                ],
            },
            {
                "head": "観察3: コスト動学は「単価90%下落 + 使用量爆発」のJevons パラドックスへ",
                "summary": [
                    "Gartner 2030年に1兆パラメータ LLM 推論コスト90%下落予測",
                    "ただしエージェント普及で総支出は増加、ROI 計算が複雑化",
                    "Sakana × NVIDIA TwELL 30%高速化など効率改善が並走",
                ],
                "cases": [
                    "Gartner「推論90%下落+総支出増」予測(2026-05-12)",
                    "Sakana × NVIDIA TwELL(2026-05-10)",
                    "Cloudflare 1,100 解雇で AI 効率化と過去最高売上が両立(2026-05-08)",
                ],
            },
            {
                "head": "観察4: 資源・電力層が「銅・水・電力」の争奪戦化",
                "summary": [
                    "Amazon が銅鉱山と直接契約、米テック大手の物資争奪が本格化",
                    "Everpure CEO「半導体メモリ4-10倍急騰」で価格70%値上げ",
                    "Stargate / SoftBank Roze / NTT データセンター譲渡で日本も巻き込まれる",
                ],
                "cases": [
                    "Amazon 銅鉱山と直接契約(2026-05-11)",
                    "Everpure 半導体4-10倍急騰(2026-05-11)",
                    "NTT データG データセンター譲渡で51%増益(2026-05-08)",
                ],
            },
        ],
        "signals_summary": [
            "AI インフラの地政学化と物理コスト爆発が同時進行",
            "宇宙データセンターという新フロンティア論議が始動",
        ],
        "signals_cases": [
            "2026-05-12  Google × SpaceX 軌道上 DC 協議",
            "2026-05-12  Google TPU 8i/8t + Willow 量子",
            "2026-05-12  Cloud Storage Rapid / クラウド Q1 35%成長",
            "2026-05-11  Amazon 銅鉱山直接契約",
        ],
        "focus_summary": [
            "NVIDIA 独占への独禁法調査(FTC/EC)の動向",
            "Jevons パラドックスの企業財務影響",
        ],
        "focus_cases": [
            "国内 AI インフラ供給(NTT/SoftBank/ファナック・安川)の主権度合い",
            "宇宙 DC の実現可能性とコスト推移",
            "半導体争奪戦下のサプライチェーン再編",
        ],
    },
    {
        "slug": "llm-evaluation",
        "title": "LLM評価・解釈可能性",
        "version": "v2",
        "articles": 7,
        "q": "ベンチマークと実用性能の乖離はどう埋まっていくか。評価指標自体はどう進化しているか。",
        "points": [
            {
                "head": "観察1: 自動ベンチマーク層は飽和、SWE-bench 型は限界に到達",
                "summary": [
                    "SWE-bench/TAU/OSWorld/GPQA Diamond/HLE などで「全社90%超」の飽和",
                    "Chatbot Arena 等の Elo 評価も差が縮小",
                    "三極化(ChatGPT/Claude/Gemini)前提の機能軸別比較が新標準",
                ],
                "cases": [
                    "ChatGPT/Claude/Gemini 三極化(2026-05-08)",
                    "llm-3services-comparison.md(2026-05-11、ユーザー作成)",
                    "@IT「日本語 AI は約1.5倍高い」(2026-05-12)",
                ],
            },
            {
                "head": "観察2: 解釈可能性層が「思考プロセス忠実度」を新評価軸に追加",
                "summary": [
                    "Natural Language Autoencoders で「答案正確性」+「思考プロセス忠実度」の2軸へ",
                    "Anthropic が解釈研究で業界標準を握る可能性",
                    "EU AI Act / 日本 AI 法律の「説明可能性」要件と直結",
                ],
                "cases": [
                    "Anthropic Natural Language Autoencoders(2026-05-07)",
                    "Anthropic 隠れ推論可視化ツール(2026-05-12)",
                    "Goodfire Silico 機械的解釈可能性公開(2026-05-01)",
                ],
            },
            {
                "head": "観察3: 専門家定性証言層がベンチマーク飽和への決定的解",
                "summary": [
                    "フィールズ賞数学者 Gowers の ChatGPT 5.5 Pro 体験記が HN 683 pts",
                    "医学(Harvard ER)、法学(Legora)などのドメイン展開が予想される",
                    "「使える/使えない」境界をトップ専門家が公開言語化",
                ],
                "cases": [
                    "Gowers の ChatGPT 5.5 Pro 体験記(2026-05-09)",
                    "Harvard 研究 AI が ER 診断で人間医師超え(2026-05-04)",
                    "AI が数学未解決問題を80分で解く(2026-05-03)",
                ],
            },
            {
                "head": "観察4: 産業応用評価が「異常検知の根拠説明」など独自軸を要求",
                "summary": [
                    "東芝 反事実波形が「AI 判断根拠の波形説明」を産業実装",
                    "NYT「Anthropic は本当に怖いか?」など主要メディアが定性評価に参入",
                    "「AI 買い物失敗5割」など消費者側の評価データも蓄積",
                ],
                "cases": [
                    "東芝 反事実波形生成技術(2026-05-12)",
                    "NYT「Anthropic の新 AI は本当に怖い?」(2026-05-12)",
                    "ITmedia 調査「AI 買い物失敗5割」(2026-05-13)",
                ],
            },
        ],
        "signals_summary": [
            "解釈可能性+定性評価+産業応用評価が同月で複数進展",
            "ピークデータ警告下で「学習データから何を学んだか」の解明が必須化",
        ],
        "signals_cases": [
            "2026-05-12  Anthropic 隠れ推論可視化ツール公開",
            "2026-05-12  東芝 反事実波形 / NYT「Anthropic 怖い?」",
            "2026-05-09  Gowers ChatGPT 5.5 Pro 体験記(HN 683 pts)",
            "2026-05-04  スタンフォード「ピークデータ」警告",
        ],
        "focus_summary": [
            "Gowers 型「専門家定性証言」の他ドメイン展開",
            "Natural Language Autoencoders の規制標準化",
        ],
        "focus_cases": [
            "日本語特化評価ベンチマークの整備",
            "「思考プロセス忠実度」のメトリクス標準化",
            "AI 監査ビジネスでの評価フレーム導入",
        ],
    },
    {
        "slug": "physical-ai",
        "title": "フィジカル AI・AI ロボティクス",
        "version": "v1",
        "articles": 0,
        "q": "フィジカル AI(ヒューマノイド/産業ロボ/自律機械)が物理空間でどんな速度で実用化されるか。VLA(視覚言語動作)モデルとハードウェアの統合は誰が主導するか。国産路線(ファナック/安川/KyoHA/Sony)vs 海外(Figure/1X/Optimus/Boston Dynamics)の競争構造はどう決着するか。",
        "points": [
            {
                "head": "観察1: ヒューマノイドロボットは「海外スタートアップ攻勢 vs 国産遅延」の構造で当面差が拡大",
                "summary": [
                    "Figure(Brett Adcock)/1X/Optimus/Apptronik/Boston Dynamics/Sanctuary AI が大型調達と量産化で先行",
                    "Meta が humanoid robotics スタートアップ買収(2026-05-02)で AI×身体性に本格参入",
                    "国内は KyoHA SEIMEI(検証機公開)、Sony フィジカル AI 布石が「研究/構想段階」、量産には遠い",
                ],
                "cases": [
                    "Meta humanoid robotics 買収(2026-05-02)",
                    "KyoHA 純国産ヒューマノイド SEIMEI 公開(2026-05-12、足首破損で動的デモ未達)",
                    "Sony CEO「TSMC 提携 + フィジカル AI 布石」(2026-05-12)",
                ],
            },
            {
                "head": "観察2: 産業ロボティクスは「国産基盤で情報資源死守」戦略が形成、SDA 主戦場化",
                "summary": [
                    "ファナック・安川電機が「フィジカル AI 国産基盤」を公式表明、情報資源の海外流出を阻止",
                    "三菱電機・オムロンが SDA(Software-Defined Automation)でデータ事業へ転換",
                    "欧米FA大手3社が「ハード売切り → ソフト継続課金」へ転換、日本勢に圧力",
                ],
                "cases": [
                    "ファナック・安川電機フィジカル AI(2026-05-11)",
                    "三菱電機・オムロン SDA でデータ事業へ(2026-05-12)",
                    "欧米FA大手の継続課金転換(2026-05-12)",
                ],
            },
            {
                "head": "観察3: VLA(視覚言語動作)モデルとハードウェア統合の「基盤レイヤ競争」が始動",
                "summary": [
                    "Nvidia Isaac が AI 基盤モデルとロボット制御の統合プラットフォームで先行",
                    "Google DeepMind co-clinician 等で「LLM × 物理アクション」の実装パターン化",
                    "国内は NEC(3D 点群90%軽量化)、PFN(PLaMo + ロボティクス)が独自路線",
                ],
                "cases": [
                    "NEC 独自 AI で 3D 点群90%軽量化(2026-05-12)",
                    "DeepMind co-clinician(医療補助 AI、2026-05-01)",
                    "尾形哲也 AI ロボット協会(AI ロボット共進化、2026-05-03)",
                ],
            },
            {
                "head": "観察4: 「データ主権 × 物理空間」が国産優位の数少ない領域として浮上",
                "summary": [
                    "ロボット動作データ・センサーデータ・現場データが国産基盤で確保される動き",
                    "国産 AI ベンダー(PFN/Sakana)との接続で「日本特化フィジカル AI」スタックの可能性",
                    "AI 主権体制(IBM Sovereign Core、日本初 AI 法律)と整合する産業政策的追い風",
                ],
                "cases": [
                    "ファナック・安川「国産基盤で情報資源死守」表明(2026-05-11)",
                    "Sony フィジカル AI 布石 + TSMC 提携(2026-05-12)",
                    "IOWN Global Forum — 通信ネットワークの主戦場は宇宙へ(2026-05-03)",
                ],
            },
        ],
        "signals_summary": [
            "ヒューマノイドロボット海外大手の量産化と国産路線の差が同月で顕在化",
            "産業ロボティクス × SDA 化が日本 FA 勢の戦略転換を迫る",
        ],
        "signals_cases": [
            "2026-05-12  KyoHA SEIMEI / Sony フィジカル AI 布石",
            "2026-05-12  三菱電機・オムロン SDA / NEC 3D 点群軽量化",
            "2026-05-11  ファナック・安川電機フィジカル AI 国産基盤",
            "2026-05-02  Meta humanoid robotics 買収",
        ],
        "focus_summary": [
            "VLA モデル × 国産 LLM(PFN/Sakana)接続の標準化タイミング",
            "ヒューマノイドロボの B2C 量産化(Figure/1X/Optimus の市場投入)",
        ],
        "focus_cases": [
            "産業ロボ × 国産 AI 基盤連携の具体プロジェクト",
            "Boston Dynamics × 韓国(Hyundai)、Apptronik × Mercedes 等の OEM 連携",
            "Sony フィジカル AI 戦略の具体プロダクト発表時期",
        ],
    },
    {
        "slug": "finance-ai",
        "title": "金融×AI(銀行・証券・保険・決済)",
        "version": "v1",
        "articles": 4,
        "q": "金融業務のどこに・どの粒度でAIが組み込まれるか。Sakana×SMBC 後の他社追随パターンと寡占化の行方。",
        "points": [
            {
                "head": "観察1: 業務適用は「営業 → アクチュアリー → 定常 → 標準化 → 寡占化」の Stage 1-2 が進行中",
                "summary": [
                    "Sakana × SMBC の提案書自動生成が Stage 1 を完了",
                    "Singular Bank × Codex で会議準備/分析が Stage 2 入り(60-90分/日削減)",
                    "MUFG × Google の商品選び〜決済で Stage 2 が消費者向けにも到達",
                ],
                "cases": [
                    "Sakana × SMBC 複数 AI エージェント提案書(2026-05-03)",
                    "Singular Bank × ChatGPT/Codex 60-90分/日削減(2026-05-06)",
                    "MUFG × Google 商品選び〜決済 AI(2026-05-07)",
                ],
            },
            {
                "head": "観察2: ベンダー競合は「Sakana × OpenAI × Anthropic × Google」の4つ巴へ",
                "summary": [
                    "Sakana(国内特化)vs PwC×OpenAI(Big4 経由)vs Anthropic Agents for FS",
                    "国内市場分割: SMBC=Sakana、MUFG=Google、Anthropic 系=Goldman/Blackstone",
                    "国産 AI ベンダー間で「業界×銀行×製品」3軸の棲み分けが進行",
                ],
                "cases": [
                    "Anthropic Agents for Financial Services(2026-05-05)",
                    "OpenAI × PwC CFO 業務再定義(2026-05-05)",
                    "Anthropic × Goldman/Blackstone $1.5B PE 所有企業向け(2026-05-05)",
                ],
            },
            {
                "head": "観察3: インフラ層は「データ基盤 × クラウド × 認証」の3課題が同時に立ち上がる",
                "summary": [
                    "三菱UFJ信託 × Snowflake、北陸銀行音声 AI などデータ基盤+業務 AI 統合",
                    "マネーフォワード GitHub 侵害で Fintech × 開発リスクが露呈",
                    "全国140信組 SKC センター8時間障害で国内金融インフラ脆弱性",
                ],
                "cases": [
                    "三菱UFJ信託 × Snowflake データマネジメント(2026-05-07)",
                    "マネーフォワード GitHub 侵害で銀行連携停止(2026-05-11)",
                    "全国140信組 SKC センター障害(2026-05-11)",
                ],
            },
            {
                "head": "観察4: 新領域として「AI 代理決済」が始動、Stripe × OpenAI が先行",
                "summary": [
                    "Stripe × OpenAI で AI が買い物する少額決済が実装段階",
                    "「AI 買い物失敗5割」など消費者側の品質問題も顕在化",
                    "Ramp $40B+ / Kalshi $22B など Fintech × AI 評価額急騰",
                ],
                "cases": [
                    "Stripe × OpenAI 少額決済(2026-05-12)",
                    "ITmedia「AI 買い物失敗5割」調査(2026-05-13)",
                    "Ramp $40B+ 評価交渉(2026-05-07)/ Kalshi $22B(2026-05-07)",
                ],
            },
        ],
        "signals_summary": [
            "国内メガバンクが米AI 直接提携で動く第一波",
            "金融 × データ基盤 × 認証の3課題が同時進行",
        ],
        "signals_cases": [
            "2026-05-13  OpenAI「How finance teams use Codex」公式",
            "2026-05-12  Stripe × OpenAI 少額決済 / Anthropic 株式二次警告",
            "2026-05-11  北陸銀行音声 AI Worker / マネフォ GitHub 侵害",
            "2026-05-07  MUFG × Google / 三菱UFJ信託 × Snowflake",
        ],
        "focus_summary": [
            "5段階モデル Stage 3(定常化)への進捗時期",
            "金融庁ガイドラインと米AI 提携の整合性",
        ],
        "focus_cases": [
            "国内 AI ベンダー間の金融市場分割の確定",
            "規制業界での Sovereign Core 採用",
            "Stripe × OpenAI 型「AI 代理経済」の国内応用",
        ],
    },
    {
        "slug": "logistics-ai",
        "title": "物流×AI(倉庫・配送・自動運転)",
        "version": "v2",
        "articles": 2,
        "q": "物流業界で AI が労働力不足・配送コスト・最適化の3課題をどう解決するか。最適化領域のホワイトカラー管理業務が AI エージェントに代替される方向性が見える中、米国先行(Amazon Robotics 等)と日本(ヤマト/佐川/SoftBank Roze 等)の戦略差はどう進化するか。",
        "points": [
            {
                "head": "観察1: AI 型分類は「ルールベース機械化 vs エージェンティック AI」で導入難易度が分岐",
                "summary": [
                    "画一化業務(ピッキング・仕分け)はルールベース機械化が現実解",
                    "動的判断(配車・例外処理)はエージェンティック AI 必須",
                    "ステップ的導入で「ルールベース → エージェント」順序が国内主流",
                ],
                "cases": [
                    "(notes.md ユーザー仮説、2026-05-11 synthesis v2 反映)",
                    "AWS × アクト・ノード見守りエージェント(2026-05-12)",
                    "Loop 物流データ基盤刷新で輸送コスト可視化(2026-05-07)",
                ],
            },
            {
                "head": "観察2: ハードウェア/ロボティクスは「国産基盤死守」戦略が形成",
                "summary": [
                    "ファナック・安川電機が「国産基盤で情報資源死守」を表明",
                    "KyoHA SEIMEI、Sony フィジカル AI 布石で国内コンソーシアム形成",
                    "中国 1X / Figure / Optimus / Sanctuary に対抗",
                ],
                "cases": [
                    "ファナック・安川電機フィジカル AI(2026-05-11)",
                    "KyoHA 純国産ヒューマノイド SEIMEI(2026-05-12)",
                    "Sony CEO TSMC 提携 + フィジカル AI 布石(2026-05-12)",
                ],
            },
            {
                "head": "観察3: データ基盤層は「SDA(Software-Defined Automation)」が主戦場",
                "summary": [
                    "欧米FA大手3社が SDA + 知能の継続課金へ転換",
                    "三菱電機・オムロンが「ハード資産+データ事業」で対抗",
                    "Loop / FourKites / project44 が物流データ基盤を整備",
                ],
                "cases": [
                    "三菱電機・オムロン SDA でデータ事業へ(2026-05-12)",
                    "欧米FA大手の継続課金転換(2026-05-12)",
                    "Loop 物流データ基盤刷新(2026-05-07)",
                ],
            },
            {
                "head": "観察4: ラストマイル層は自動運転調整局面 + 一次産業エージェント",
                "summary": [
                    "Kodiak AI $100M ディスカウント調達で自動運転トラック業界調整",
                    "Rivian AI 音声アシスタント全車展開でEVに AI が標準化",
                    "AWS × アクト・ノード「見守りエージェント」で農林漁業 AI 化",
                ],
                "cases": [
                    "Kodiak AI $100M ディスカウント調達 株価37%下落(2026-05-07)",
                    "Rivian AI 音声アシスタント全車展開(2026-05-12)",
                    "AWS × アクト・ノード見守りエージェント 工数50%削減(2026-05-12)",
                ],
            },
        ],
        "signals_summary": [
            "国産フィジカル AI コンソーシアムが同月で複数表面化",
            "欧米 FA 大手の SDA × 継続課金転換が日本に波及",
        ],
        "signals_cases": [
            "2026-05-12  KyoHA SEIMEI / Sony フィジカル AI",
            "2026-05-12  三菱電機・オムロン SDA",
            "2026-05-11  ファナック・安川電機フィジカル AI",
            "2026-05-07  Loop 物流データ基盤刷新",
        ],
        "focus_summary": [
            "VLA モデル × 国産 LLM(PFN/Sakana)接続の標準化",
            "国産フィジカル AI コンソーシアムの拡大",
        ],
        "focus_cases": [
            "ヤマト / 佐川 / 日本郵便の AI 採用パターン",
            "自動運転トラックの北米 vs 国内戦略差",
            "SDA × 国内 FA × 国産 LLM の3軸連携",
        ],
    },
    # NOTE:  トピックは公開停止のため削除済み
    {
        "slug": "consulting-ai",
        "title": "コンサルティング業界×AI",
        "version": "v1",
        "articles": 11,
        "q": "コンサルティング業界(Big4/戦略系/SIer/シンクタンク)で AI が競争優位・収益モデル・組織構造をどう書き換えるか。作業の AI 化、提供価値再定義、3層分化、日本市場の特性、シンクタンク × 自治体、AI ネイティブの脅威 — を継続観察する。",
        "points": [
            {
                "head": "観察1: プレイヤー位置取りは「5層分化」で確定しつつある",
                "summary": [
                    "戦略立案(McKinsey/BCG/Bain)/ 変革管理(NRI・Big4)",
                    "実装 SI(Accenture/IBM/Capgemini/NEC・富士通・日立)",
                    "公共規制(Big4 監査 + シンクタンク)/ AI ネイティブ侵入(OpenAI/Anthropic/Sierra/Cognition)",
                ],
                "cases": [
                    "OpenAI Deployment Co + $4B + Tomoro(2026-05-12)",
                    "NRI「AI と組織の未来」(2026-05-12)",
                    "Cognition AI 日本法人設立(2026-04-29)",
                ],
            },
            {
                "head": "観察2: 提供価値は「実装層が AI ベンダー直営 FDE モデルで侵食」される",
                "summary": [
                    "OpenAI Deployment Co + Tomoro 買収で FDE モデルが汎用 AI 業界に到来",
                    "Palantir Forward Deployed Engineer モデルが他社にも拡散",
                    "コンサル中抜き(中抜き型) vs コンサル取り込み(取り込み型)の両方が並走",
                ],
                "cases": [
                    "OpenAI Deployment Co 設立 + 19社設立パートナー(2026-05-12)",
                    "アクセンチュア × Anthropic 国内本格化(2026-05-12)",
                    "OpenAI × PwC CFO 業務再定義(2026-05-05)",
                ],
            },
            {
                "head": "観察3: 国内市場は「人月脱却 / 規模拡大 / 米AI 提携」の3戦略に分岐",
                "summary": [
                    "富士ソフト(成果報酬移行)/ 日立 GlobalLogic(規模拡大)/ NEC × Anthropic(米AI 提携)",
                    "中堅 SI(TIS/SCSK/BIPROGY/CTC)の戦略選択が次の試金石",
                    "シンクタンク(NRI/三菱UFJリサーチ)が「公共セクター実プロダクト化」に進出",
                ],
                "cases": [
                    "富士ソフト SIer 新体制 人月商売脱却(2026-04-30)",
                    "日立 GlobalLogic + Hitachi Digital Services 統合(2026-05-11)",
                    "三菱UFJリサーチ Gov Sales(2026-05-12)",
                ],
            },
            {
                "head": "観察4: AI ネイティブ脅威は「Tier 2-3 領域を吸収」段階に到達",
                "summary": [
                    "Cognition AI(Devin)が自律 AI エンジニア時間貸しモデルを国内展開",
                    "Sierra AI、Decart、Lovable などが SaaS / SI / コンサル を構造的に置換",
                    "戦略系の上位フィー圧迫(GenAI 戦略コンサルの単価半減リスク)",
                ],
                "cases": [
                    "Cognition AI 日本法人設立(2026-04-29)",
                    "ダイヤモンド「AI ネイティブスタートアップの脅威」(2026-05-11)",
                    "OpenAI How finance teams use Codex(2026-05-12、Big4 直接競合の続編)",
                ],
            },
        ],
        "signals_summary": [
            "AI ベンダーが業界・コンサル両軸で直接侵入を加速",
            "国内 SI の戦略分岐が同月で複数表面化",
        ],
        "signals_cases": [
            "2026-05-12  アクセンチュア × Anthropic 国内本格化 + IPA/デロイト不正",
            "2026-05-12  OpenAI Deployment Co + Tomoro / Anthropic Claude for Legal",
            "2026-05-11  NEC × Anthropic / 日立 GlobalLogic 統合",
            "2026-05-04  富士ソフト SIer 新体制 人月脱却",
        ],
        "focus_summary": [
            "Big4 全社の AI 提携完了タイミング(Deloitte/EY/KPMG)",
            "国内 SI 決算で人月モデル崩壊度合いが定量化されるか",
        ],
        "focus_cases": [
            "OpenAI Deployment Co の初期成果がコンサル収益にどう波及",
            "AI 監査ビジネス(Big4)の市場規模立ち上がり",
            "Sakana・PFN 等が同様の米AI 提携を選ぶか独自路線継続か",
        ],
    },
]

# ---------- Slide 1: combined overview + daily ----------
def build_combined_slide(prs, page_num, page_total):
    s = blank_slide(prs)
    # ヘッダ
    add_rect(s, Inches(0.0), Inches(0.0), Inches(0.18), SLIDE_H, fill=SAGE)
    add_text(s, Inches(0.55), Inches(0.32), Inches(11.5), Inches(0.32),
             "AI NEWS VAULT  ·  TOPIC OVERVIEW  +  DAILY 2026-05-13",
             size=10, bold=True, color=SAGE, font=FONT_JP)
    add_text(s, Inches(0.55), Inches(0.62), Inches(11.5), Inches(0.7),
             "AI業界の潮流 × 本日のハイライト",
             size=24, bold=True, color=DEEP_INDIGO, font=FONT_JP)
    add_text(s, Inches(12.0), Inches(0.32), Inches(1.0), Inches(0.32),
             f"{page_num:02d} / {page_total:02d}", size=10, bold=True,
             color=SOFT_GRAY, font=FONT_JP, align=PP_ALIGN.RIGHT)
    add_line(s, Inches(0.55), Inches(1.4), Inches(12.85), Inches(1.4),
             color=STONE_MID, weight_pt=0.75)
    add_line(s, Inches(6.665), Inches(1.55), Inches(6.665), Inches(7.05),
             color=STONE_MID, weight_pt=0.5)

    # ====== 左半分: 3つの大潮流 ======
    L_x = Inches(0.55)
    L_w = Inches(5.85)
    add_text(s, L_x, Inches(1.55), L_w, Inches(0.32),
             "■ 13トピック横断 — AI業界の3つの大潮流",
             size=11, bold=True, color=SAGE, font=FONT_JP)

    # --- 潮流01: 3軸整理 ---
    y01 = Inches(1.95)
    h01 = Inches(1.8)
    add_rect(s, L_x, y01, L_w, h01, fill=STONE_LIGHT)
    add_text(s, L_x + Inches(0.2), y01 + Inches(0.12), Inches(0.7), Inches(0.45),
             "01", size=22, bold=True, color=SAGE, font=FONT_JP)
    add_text(s, L_x + Inches(0.95), y01 + Inches(0.1), L_w - Inches(1.1),
             Inches(0.32), "コンサル×AIベンダー×業界Tier1 — 3軸の同時進行",
             size=12, bold=True, color=DEEP_INDIGO, font=FONT_JP)
    # 概要
    add_text(s, L_x + Inches(0.95), y01 + Inches(0.42), L_w - Inches(1.1),
             Inches(0.3),
             "概要: AIベンダーがコンサル業界に「取り込み(a)」「中抜き(b)」「業界直結(c)」の3軸で侵入",
             size=8, bold=True, color=SLATE, font=FONT_JP)
    # 3軸事例
    axes = [
        ("(a) AIベンダー×コンサル", "アクセンチュア×Anthropic 国内 / NEC×Anthropic / OpenAI Deployment Co + Tomoro(中抜き型)"),
        ("(b) AIベンダー×業界Tier1", "Anthropic Claude for Legal × Big Law / OpenAI×PwC / Sakana×SMBC / Stripe×OpenAI"),
        ("(c) コンサル×業界Tier1", "アクセンチュア×SAP ADVANCE / アクセンチュア×日本精工 / 富士ソフト人月脱却"),
    ]
    for i, (label, body) in enumerate(axes):
        yy = y01 + Inches(0.75 + 0.32 * i)
        add_text(s, L_x + Inches(0.95), yy, Inches(2.0), Inches(0.28),
                 label, size=8, bold=True, color=TERRACOTTA, font=FONT_JP)
        add_text(s, L_x + Inches(2.95), yy, L_w - Inches(3.1), Inches(0.28),
                 body, size=8, color=SLATE, font=FONT_JP)

    # --- 潮流02: 業務代替×AIネイティブ ---
    y02 = y01 + h01 + Inches(0.1)
    h02 = Inches(1.6)
    add_rect(s, L_x, y02, L_w, h02, fill=STONE_LIGHT)
    add_text(s, L_x + Inches(0.2), y02 + Inches(0.12), Inches(0.7), Inches(0.45),
             "02", size=22, bold=True, color=SAGE, font=FONT_JP)
    add_text(s, L_x + Inches(0.95), y02 + Inches(0.1), L_w - Inches(1.1),
             Inches(0.32),
             "業務代替型AI解雇とAIネイティブ企業の台頭",
             size=12, bold=True, color=DEEP_INDIGO, font=FONT_JP)
    add_text(s, L_x + Inches(0.95), y02 + Inches(0.42), L_w - Inches(1.1),
             Inches(0.55),
             "概要: 「AI効率化型解雇」がCEO公言レベルに達し、AIネイティブ企業が既存SaaS/SI/コンサルを構造的に置換する圧力が顕在化",
             size=8, bold=True, color=SLATE, font=FONT_JP)
    add_text(s, L_x + Inches(0.95), y02 + Inches(0.92), L_w - Inches(1.1),
             Inches(0.6),
             "事例: Cloudflare 1,100 / Coinbase 14% / GM 数百名 / GitLab / "
             "Gartner 28% CEO「AI 収益リスク」/ ダイヤモンド「AIネイティブスタートアップ脅威」/ "
             "Cognition AI日本法人 / Sierra AI / @IT「AI 業務浸食予想以上」",
             size=8, color=SLATE, font=FONT_JP)

    # --- 潮流03: 国内 AI 主権体制 ---
    y03 = y02 + h02 + Inches(0.1)
    h03 = Inches(1.7)
    add_rect(s, L_x, y03, L_w, h03, fill=STONE_LIGHT)
    add_text(s, L_x + Inches(0.2), y03 + Inches(0.12), Inches(0.7), Inches(0.45),
             "03", size=22, bold=True, color=SAGE, font=FONT_JP)
    add_text(s, L_x + Inches(0.95), y03 + Inches(0.1), L_w - Inches(1.1),
             Inches(0.32),
             "国内AI主権体制の同時成立(2026-05-11週)",
             size=12, bold=True, color=DEEP_INDIGO, font=FONT_JP)
    # 概要 箇条書き
    add_text(s, L_x + Inches(0.95), y03 + Inches(0.42), L_w - Inches(1.1),
             Inches(0.28),
             "▸ 法律/基盤製品/米AI提携/製造業適応の4層が単一週で同時表面化",
             size=8, color=SLATE, font=FONT_JP)
    add_text(s, L_x + Inches(0.95), y03 + Inches(0.65), L_w - Inches(1.1),
             Inches(0.28),
             "▸ EU規制重視・米自由市場とは別の「公共財型 AI 戦略」が形成",
             size=8, color=SLATE, font=FONT_JP)
    # 事例 4層
    detail = [
        ("法律", "日本初 AI 特化法律(2026-05-11) / サイバー対処強化法10月施行"),
        ("基盤製品", "IBM Sovereign Core 一般提供(2026-05-05/11) / IBM Enterprise Advantage 拡充"),
        ("米AI提携", "NEC×Anthropic / アクセンチュア×Anthropic / Claude Platform on AWS"),
        ("製造業適応", "ファナック・安川フィジカルAI / 三菱電機・オムロン SDA SAT"),
    ]
    for i, (label, body) in enumerate(detail):
        yy = y03 + Inches(0.95 + 0.18 * i)
        add_text(s, L_x + Inches(0.95), yy, Inches(1.2), Inches(0.18),
                 label, size=7, bold=True, color=TERRACOTTA, font=FONT_JP)
        add_text(s, L_x + Inches(2.15), yy, L_w - Inches(2.3), Inches(0.18),
                 body, size=7, color=SLATE, font=FONT_JP)

    # ====== 右半分: 本日のDaily ======
    R_x = Inches(6.85)
    R_w = Inches(6.0)
    add_text(s, R_x, Inches(1.55), R_w, Inches(0.32),
             "■ 本日のハイライト  2026-05-13(93記事中抜粋)",
             size=11, bold=True, color=TERRACOTTA, font=FONT_JP)

    pill_y = Inches(1.95)
    themes = [
        ("法務AI", SAGE),
        ("国内SI提携", SAGE),
        ("Musk v Altman", TERRACOTTA),
        ("インフラ宇宙化", SAGE),
        ("AI兵器化", TERRACOTTA),
    ]
    px = R_x
    for label, color in themes:
        w = Inches(0.9 + 0.1 * (len(label) - 3))
        add_pill(s, px, pill_y, w, Inches(0.28), label,
                 fill=color, size=9)
        px += w + Inches(0.08)

    items = [
        ("Anthropic Claude for Legal 一斉ローンチ",
         "20+ コネクタ / 12 プラクティスエリアプラグイン / Thomson Reuters CoCounsel 連携。Big Law が AI に all-in。"),
        ("アクセンチュア × Anthropic 国内本格化",
         "Claude の4支援領域。Claude Platform on AWS 正式リリースと同日 → Anthropic 国内3軸提携完成。"),
        ("Sam Altman 公判で証言開始",
         "「Musk が huge damage」「Musk は OpenAI を子供たちに渡そうとした」と証言。"),
        ("Google × SpaceX、軌道上 AI データセンターで協議",
         "AI 計算インフラの宇宙化。SpaceX Terafab、Amazon 銅鉱山と並ぶ次段階。"),
        ("Google Android Show — Googlebooks + Gemini Intelligence",
         "AI-first ChromeBook / vibe-coded widgets / Gemini in Chrome / Gboard 音声入力。"),
        ("ClaudeBleed 脆弱性発覚",
         "Claude for Chrome 拡張で AI 操作乗っ取り可能(LayerX Security)。修正後も残存。"),
        ("AI 製ゼロデイ + 中朝露国家関与",
         "Google 公式報告。HackerOne / Mythos 専門家250人 / 高市総理サイバー指示と並走。"),
        ("HEROZ ASK が MCP 対応開始",
         "国内法人向け SaaS の MCP 第一号。Notion / Box / Slack 連携。"),
    ]
    li_y = Inches(2.4)
    row_h = Inches(0.56)
    for i, (head, body) in enumerate(items):
        y_i = li_y + row_h * i
        add_rect(s, R_x, y_i + Inches(0.06), Inches(0.32), Inches(0.32),
                 fill=DEEP_INDIGO)
        add_text(s, R_x, y_i + Inches(0.08), Inches(0.32), Inches(0.28),
                 f"{i+1:02d}", size=9, bold=True, color=WHITE,
                 font=FONT_JP, align=PP_ALIGN.CENTER)
        add_text(s, R_x + Inches(0.42), y_i + Inches(0.02),
                 R_w - Inches(0.42), Inches(0.26),
                 head, size=10, bold=True, color=DEEP_INDIGO, font=FONT_JP)
        add_text(s, R_x + Inches(0.42), y_i + Inches(0.27),
                 R_w - Inches(0.42), Inches(0.3),
                 body, size=8, color=SLATE, font=FONT_JP)

    footer(s)

# ---------- Topic slide (new layout) ----------
def build_topic_slide(prs, topic, page_num, page_total):
    s = blank_slide(prs)
    eyebrow = f"TOPIC  ·  {topic['slug']}  ·  synthesis {topic['version']}  ·  articles: {topic['articles']}"
    header_band(s, eyebrow, topic["title"], page_num, page_total)

    # 自分の問い(横長)
    q_y = Inches(1.5)
    add_text(s, Inches(0.55), q_y, Inches(1.4), Inches(0.28),
             "自分の問い", size=10, bold=True, color=SAGE, font=FONT_JP)
    add_text(s, Inches(0.55), q_y + Inches(0.28), Inches(12.3), Inches(0.4),
             topic["q"], size=11, color=DEEP_INDIGO, font=FONT_JP)

    # 2-column layout
    L_x = Inches(0.55)
    L_w = Inches(7.6)
    R_x = Inches(8.35)
    R_w = Inches(4.7)
    content_y = Inches(2.25)

    # ===== 左: 主要論点 =====
    add_text(s, L_x, content_y, L_w, Inches(0.3),
             "■ 主要論点(問い検証の観察+含意)", size=10, bold=True,
             color=SAGE, font=FONT_JP)
    points = topic["points"]
    n = len(points)
    pt_top = content_y + Inches(0.35)
    pt_total = Inches(4.6)
    pt_h = pt_total / n
    for i, pt in enumerate(points):
        py = pt_top + pt_h * i
        # 番号
        add_text(s, L_x, py, Inches(0.4), Inches(0.3),
                 f"{i+1:02d}", size=14, bold=True, color=TERRACOTTA, font=FONT_JP)
        # 観察+含意
        add_text(s, L_x + Inches(0.42), py - Inches(0.02), L_w - Inches(0.42),
                 Inches(0.4), pt["head"], size=10, bold=True,
                 color=DEEP_INDIGO, font=FONT_JP)
        # 概要 + 事例(コンパクト箇条書き)
        body_y = py + Inches(0.32)
        # 概要 label
        add_text(s, L_x + Inches(0.42), body_y, Inches(0.7), Inches(0.18),
                 "概要", size=7, bold=True, color=SAGE, font=FONT_JP)
        # 概要 items
        sm_x = L_x + Inches(0.42)
        for j, item in enumerate(pt["summary"]):
            yy = body_y + Inches(0.18 + 0.18 * j)
            add_text(s, sm_x, yy, Inches(0.15), Inches(0.16),
                     "・", size=7, color=SAGE, font=FONT_JP)
            add_text(s, sm_x + Inches(0.13), yy, L_w - Inches(0.55),
                     Inches(0.18), item, size=7, color=SLATE, font=FONT_JP)
        # 事例 label
        cases_start = body_y + Inches(0.18 + 0.18 * len(pt["summary"]) + 0.04)
        add_text(s, L_x + Inches(0.42), cases_start, Inches(0.7), Inches(0.18),
                 "事例", size=7, bold=True, color=TERRACOTTA, font=FONT_JP)
        for j, item in enumerate(pt["cases"]):
            yy = cases_start + Inches(0.18 + 0.18 * j)
            add_text(s, sm_x, yy, Inches(0.15), Inches(0.16),
                     "▸", size=7, color=TERRACOTTA, font=FONT_JP)
            add_text(s, sm_x + Inches(0.13), yy, L_w - Inches(0.55),
                     Inches(0.18), item, size=7, color=SLATE, font=FONT_JP)

    # ===== 右: シグナル + 注視点 =====
    # シグナル
    add_text(s, R_x, content_y, R_w, Inches(0.3),
             "■ 直近の重要シグナル(2026-04〜05)", size=10, bold=True,
             color=SAGE, font=FONT_JP)
    sig_top = content_y + Inches(0.35)
    add_rect(s, R_x, sig_top, R_w, Inches(2.55), fill=STONE_LIGHT)

    # 概要
    add_text(s, R_x + Inches(0.15), sig_top + Inches(0.12),
             Inches(0.7), Inches(0.18),
             "概要", size=7, bold=True, color=SAGE, font=FONT_JP)
    for j, item in enumerate(topic["signals_summary"]):
        yy = sig_top + Inches(0.32 + 0.22 * j)
        add_text(s, R_x + Inches(0.18), yy, Inches(0.15), Inches(0.2),
                 "・", size=8, color=SAGE, font=FONT_JP)
        add_text(s, R_x + Inches(0.32), yy, R_w - Inches(0.45),
                 Inches(0.22), item, size=8, color=SLATE, font=FONT_JP)
    # 事例 label
    sig_cases_y = sig_top + Inches(0.32 + 0.22 * len(topic["signals_summary"]) + 0.05)
    add_text(s, R_x + Inches(0.15), sig_cases_y, Inches(0.7), Inches(0.18),
             "事例", size=7, bold=True, color=TERRACOTTA, font=FONT_JP)
    for j, item in enumerate(topic["signals_cases"]):
        yy = sig_cases_y + Inches(0.22 + 0.22 * j)
        add_text(s, R_x + Inches(0.18), yy, Inches(0.15), Inches(0.2),
                 "▸", size=8, color=TERRACOTTA, font=FONT_JP)
        add_text(s, R_x + Inches(0.32), yy, R_w - Inches(0.45),
                 Inches(0.22), item, size=8, color=SLATE, font=FONT_JP)

    # 注視点
    focus_top = sig_top + Inches(2.65)
    add_text(s, R_x, focus_top, R_w, Inches(0.3),
             "■ 次の注視点", size=10, bold=True,
             color=TERRACOTTA, font=FONT_JP)
    focus_box_top = focus_top + Inches(0.35)
    add_rect(s, R_x, focus_box_top, R_w, Inches(1.95), fill=STONE_LIGHT)
    add_text(s, R_x + Inches(0.15), focus_box_top + Inches(0.12),
             Inches(0.7), Inches(0.18),
             "概要", size=7, bold=True, color=SAGE, font=FONT_JP)
    for j, item in enumerate(topic["focus_summary"]):
        yy = focus_box_top + Inches(0.32 + 0.22 * j)
        add_text(s, R_x + Inches(0.18), yy, Inches(0.15), Inches(0.2),
                 "・", size=8, color=SAGE, font=FONT_JP)
        add_text(s, R_x + Inches(0.32), yy, R_w - Inches(0.45),
                 Inches(0.22), item, size=8, color=SLATE, font=FONT_JP)
    fo_cases_y = focus_box_top + Inches(0.32 + 0.22 * len(topic["focus_summary"]) + 0.05)
    add_text(s, R_x + Inches(0.15), fo_cases_y, Inches(0.7), Inches(0.18),
             "事例/論点", size=7, bold=True, color=TERRACOTTA, font=FONT_JP)
    for j, item in enumerate(topic["focus_cases"]):
        yy = fo_cases_y + Inches(0.22 + 0.22 * j)
        add_text(s, R_x + Inches(0.18), yy, Inches(0.15), Inches(0.2),
                 "▸", size=8, color=TERRACOTTA, font=FONT_JP)
        add_text(s, R_x + Inches(0.32), yy, R_w - Inches(0.45),
                 Inches(0.22), item, size=8, color=SLATE, font=FONT_JP)

    footer(s)

# ----------  Impact Slide ----------
def main():
    root = Path(__file__).resolve().parent.parent
    out_dir = root / "topics" / "_meta"
    out_dir.mkdir(parents=True, exist_ok=True)
    out_path = out_dir / "topics-overview-2026-05-13-v3.pptx"

    prs = new_deck()
    page_total = 1 + len(TOPICS) + 1  # combined + topics +  impact
    build_combined_slide(prs, page_num=1, page_total=page_total)
    for i, topic in enumerate(TOPICS):
        build_topic_slide(prs, topic, page_num=2 + i, page_total=page_total)
    prs.save(str(out_path))
    print(f"Saved: {out_path}")

    docs = Path.home() / "Documents"
    if docs.exists():
        dest = docs / "ai-news-vault_topics-overview-2026-05-13-v3.pptx"
        shutil.copyfile(out_path, dest)
        print(f"Copied: {dest}")

if __name__ == "__main__":
    main()
